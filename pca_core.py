# src/pca_core.py

import os
from io import BytesIO
from typing import Optional, List

import pandas as pd

import matplotlib
matplotlib.use("Agg")  # safe for servers / headless envs
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from dotenv import load_dotenv
from openai import OpenAI


# -------------------------------------------------
# Load environment & OpenAI client
# -------------------------------------------------
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    print("[WARN] OPENAI_API_KEY is not set. LLM calls will fail.")

client: Optional[OpenAI] = None
if OPENAI_API_KEY:
    client = OpenAI(api_key=OPENAI_API_KEY)


# -------------------------------------------------
# Theme colors & helpers
# -------------------------------------------------
THEME_BG = RGBColor(249, 250, 251)    # soft light background
TITLE_COLOR = RGBColor(15, 23, 42)    # dark navy
TEXT_COLOR = RGBColor(55, 65, 81)     # gray
ACCENT = RGBColor(79, 70, 229)        # indigo
ACCENT_SOFT = RGBColor(165, 180, 252) # light indigo


def _style_slide_background(slide):
    """Apply a light theme background to the slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME_BG


def _style_title(shape):
    """Style the title shape with theme color & size."""
    if shape is None:
        return
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = p.text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def _add_bullets(placeholder, bullets: List[str], max_bullets: int, font_size: int = 20):
    """Add up to max_bullets themed bullets to a text placeholder."""
    tf = placeholder.text_frame
    tf.clear()
    for b in bullets[:max_bullets]:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.level = 0
        if p.runs:
            r = p.runs[0]
        else:
            r = p.add_run()
        r.font.size = Pt(font_size)
        r.font.color.rgb = TEXT_COLOR


# =================================================
# 1. DATA SUMMARY – short, only for LLM context
# =================================================
def _summarise_dataframe(df: pd.DataFrame) -> str:
    """
    Very short, high-level summary of the dataset, only for LLM context.
    No raw rows / full stats.
    """
    if df.empty:
        return "No structured campaign data uploaded."

    lines: list[str] = []
    lines.append(f"Dataset shape: {df.shape[0]} rows x {df.shape[1]} columns.")

    # Channel / dimension columns
    for col in ["Channel", "Platform", "Publisher"]:
        if col in df.columns:
            top = df[col].value_counts().head(3)
            vals = ", ".join(f"{k} ({v})" for k, v in top.items())
            lines.append(f"{col}: {vals}")

    # Numeric KPIs snapshot
    num = df.select_dtypes(include="number")
    if not num.empty:
        lines.append("Numeric KPIs (sample):")
        for col in num.columns[:4]:
            lines.append(
                f"- {col}: mean={num[col].mean():.2f}, "
                f"min={num[col].min():.2f}, max={num[col].max():.2f}"
            )

    return "\n".join(lines)


# =================================================
# 2. LLM HELPERS – short bullet outputs
# =================================================
def _ensure_client() -> OpenAI:
    if client is None:
        raise RuntimeError(
            "OPENAI_API_KEY is missing or invalid. "
            "Set it in your .env file as OPENAI_API_KEY=sk-..."
        )
    return client


def _generate_llm_insights(
    objective: str,
    data_summary: str,
    channels: List[str],
    has_image: bool,
) -> str:
    """
    Generate short, bullet-point insights for slides.
    """
    c = _ensure_client()

    system_msg = (
        "You are a senior digital marketing analyst. "
        "You write ONLY short bullet point insights suitable for PPT slides. "
        "No markdown headings, no long paragraphs."
    )

    user_msg = f"""
CAMPAIGN OBJECTIVE:
{objective or '(no explicit objective text provided)'}

CHANNELS:
{", ".join(channels) if channels else "Not specified"}

DATA SNAPSHOT (for your context, do NOT repeat numbers verbatim):
{data_summary}

DASHBOARD SNAPSHOT PROVIDED: {"Yes" if has_image else "No"}.

TASK:
Write 8–12 bullet points that cover:
- what worked well
- what underperformed
- KPI behaviour
- big insights
- key risks / watchouts

RULES:
- Each bullet max ~20 words.
- Do NOT use markdown ### or numbered sections.
- Start each line with a bullet like: • text...
"""

    try:
        resp = c.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg},
            ],
            temperature=0.4,
            max_tokens=1000,
        )
        return resp.choices[0].message.content
    except Exception as e:
        print("[ERROR] LLM in _generate_llm_insights:", repr(e))
        return f"• (LLM insights failed: {e})"


def _generate_llm_recos(objective: str, insights_text: str) -> str:
    """
    Generate future-facing recommendations and next steps as bullets.
    """
    c = _ensure_client()

    prompt = f"""
Objective:
{objective or '(no explicit objective provided)'}

Insights:
{insights_text}

TASK:
Write 5–8 future-facing recommendations and next steps.
Each bullet:
- one clear action
- max ~20 words
- start with: •

No markdown headings, no numbering, only bullet lines.
"""

    try:
        resp = c.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a performance marketing strategist."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.4,
            max_tokens=600,
        )
        return resp.choices[0].message.content
    except Exception as e:
        print("[ERROR] LLM in _generate_llm_recos:", repr(e))
        return f"• (LLM recommendations failed: {e})"


# =================================================
# 3. CHART HELPERS – visuals from data
# =================================================
def _overall_trend_chart(df: pd.DataFrame) -> Optional[BytesIO]:
    """
    Simple trend chart based on the first numeric column.
    """
    num = df.select_dtypes(include="number")
    if num.empty:
        return None

    col = num.columns[0]
    plt.figure(figsize=(7, 3))
    plt.plot(num[col].values[:50], marker="o")
    plt.title(f"Trend – {col}")
    plt.xlabel("Index")
    plt.ylabel(col)
    plt.grid(True, alpha=0.3)
    plt.tight_layout()

    buf = BytesIO()
    plt.savefig(buf, format="png")
    plt.close()
    buf.seek(0)
    return buf


def _channel_perf_chart(df: pd.DataFrame) -> Optional[BytesIO]:
    """
    Channel performance bar chart if a column like 'Channel' exists.
    Uses the first numeric metric.
    """
    channel_col = None
    for candidate in ["Channel", "channel", "Platform", "platform"]:
        if candidate in df.columns:
            channel_col = candidate
            break

    if channel_col is None:
        return None

    num = df.select_dtypes(include="number")
    if num.empty:
        return None

    metric_col = num.columns[0]
    grouped = df.groupby(channel_col)[metric_col].sum().sort_values(ascending=False).head(6)

    plt.figure(figsize=(7, 3))
    grouped.plot(kind="bar")
    plt.title(f"{metric_col} by {channel_col}")
    plt.ylabel(metric_col)
    plt.xticks(rotation=20, ha="right")
    plt.grid(axis="y", alpha=0.3)
    plt.tight_layout()

    buf = BytesIO()
    plt.savefig(buf, format="png")
    plt.close()
    buf.seek(0)
    return buf


# =================================================
# 4. BULLET NORMALIZATION
# =================================================
def _extract_bullets(raw: str) -> List[str]:
    """
    Normalize raw LLM text to a clean list of bullet strings.
    """
    bullets: List[str] = []
    for line in raw.splitlines():
        line = line.strip()
        if not line:
            continue
        # strip common bullet markers
        for prefix in ("•", "-", "*", "–", "—"):
            if line.startswith(prefix):
                line = line[len(prefix):].strip()
        if not line or line.startswith("#"):
            continue
        bullets.append(line)
    return bullets


# =================================================
# 5. PPT BUILDER – THEMED, VISUAL, CLIENT-READY
# =================================================
def build_final_ppt(
    insights_text: str,
    recos_text: str,
    df: pd.DataFrame,
    image_bytes: Optional[bytes] = None,
) -> BytesIO:
    """
    Generates a themed, client-ready PPT deck:
      1. Title
      2. Executive Summary
      3. Key Insights
      4. KPI Overview (table)
      5. Performance Trend (chart)
      6. Channel Performance (chart)
      7. Recommendations
      8. Next Steps
      + optional dashboard snapshot if image exists
    """
    prs = Presentation()

    insight_bullets = _extract_bullets(insights_text)
    reco_bullets = _extract_bullets(recos_text)

    # -------- Slide 1: Title --------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _style_slide_background(slide)
    slide.shapes.title.text = "Post Campaign Analysis"
    _style_title(slide.shapes.title)

    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = "Automated, AI-driven insights & recommendations"
        p = sub.text_frame.paragraphs[0]
        if p.runs:
            r = p.runs[0]
        else:
            r = p.add_run()
        r.font.size = Pt(20)
        r.font.color.rgb = TEXT_COLOR

    # -------- Slide 2: Executive Summary --------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _style_slide_background(slide)
    slide.shapes.title.text = "Executive Summary"
    _style_title(slide.shapes.title)

    top_exec = insight_bullets[:5] if insight_bullets else ["No insights generated."]
    _add_bullets(slide.placeholders[1], top_exec, max_bullets=5, font_size=20)

    # -------- Slide 3: Key Insights --------
    if len(insight_bullets) > 5:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        _style_slide_background(slide)
        slide.shapes.title.text = "Key Insights"
        _style_title(slide.shapes.title)

        remaining = insight_bullets[5:]
        _add_bullets(slide.placeholders[1], remaining, max_bullets=8, font_size=18)

    # -------- Slide 4: KPI Overview (table from numeric cols) --------
    if not df.empty:
        numeric = df.select_dtypes(include="number")
        if not numeric.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            _style_slide_background(slide)
            slide.shapes.title.text = "KPI Overview"
            _style_title(slide.shapes.title)

            # Use first row as "representative" snapshot
            first_row = numeric.head(1).iloc[0]
            metrics = list(first_row.index[:6])
            values = [first_row[m] for m in metrics]

            rows = len(metrics) + 1
            cols = 2
            left = Inches(1)
            top = Inches(1.7)
            width = Inches(8)
            height = Inches(2.5)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            table.columns[0].width = Inches(3.5)
            table.columns[1].width = Inches(2.5)

            # Header row
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Value"

            for i, (metric, val) in enumerate(zip(metrics, values), start=1):
                table.cell(i, 0).text = str(metric)
                try:
                    table.cell(i, 1).text = f"{float(val):.2f}"
                except Exception:
                    table.cell(i, 1).text = str(val)

    # -------- Slide 5: Performance Trend (chart) --------
    if not df.empty:
        trend_buf = _overall_trend_chart(df)
        if trend_buf:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            _style_slide_background(slide)
            slide.shapes.title.text = "Performance Trend"
            _style_title(slide.shapes.title)

            left = Inches(1)
            top = Inches(1.7)
            slide.shapes.add_picture(trend_buf, left, top, width=Inches(8))

    # -------- Slide 6: Channel Performance (chart) --------
    if not df.empty:
        ch_buf = _channel_perf_chart(df)
        if ch_buf:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            _style_slide_background(slide)
            slide.shapes.title.text = "Channel Performance"
            _style_title(slide.shapes.title)

            left = Inches(1)
            top = Inches(1.7)
            slide.shapes.add_picture(ch_buf, left, top, width=Inches(8))

    # -------- Optional Dashboard Snapshot (you can treat as Slide 7) --------
    if image_bytes:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _style_slide_background(slide)
        slide.shapes.title.text = "Dashboard Snapshot"
        _style_title(slide.shapes.title)

        left = Inches(0.7)
        top = Inches(1.7)
        slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=Inches(8.5))

    # -------- Recommendations --------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _style_slide_background(slide)
    slide.shapes.title.text = "Recommendations"
    _style_title(slide.shapes.title)

    if not reco_bullets:
        reco_bullets = ["No recommendations generated due to an error."]
    _add_bullets(slide.placeholders[1], reco_bullets, max_bullets=7, font_size=18)

    # -------- Next Steps --------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _style_slide_background(slide)
    slide.shapes.title.text = "Next Steps"
    _style_title(slide.shapes.title)

    next_steps = [
        "Review these insights with client stakeholders.",
        "Agree on priority tests and optimizations for the next cycle.",
        "Translate recommendations into a concrete experimentation roadmap.",
        "Refresh creatives and messaging based on top-performing segments.",
        "Refine tracking & attribution to reduce data blind spots.",
    ]
    _add_bullets(slide.placeholders[1], next_steps, max_bullets=6, font_size=18)

    # Save to memory
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# =================================================
# 6. ENTRYPOINT – USED BY FastAPI main.py
# =================================================
def run_pca_pipeline(
    image_bytes: Optional[bytes] = None,
    data_bytes: Optional[bytes] = None,
    data_filename: Optional[str] = None,
    objective: str = "",
    channels: Optional[list] = None,
) -> BytesIO:
    """
    Main function used by FastAPI endpoint.

    Parameters
    ----------
    image_bytes : bytes | None
        Raw bytes of uploaded dashboard image (PNG/JPG).
    data_bytes : bytes | None
        Raw bytes of uploaded CSV/XLSX.
    data_filename : str | None
        Name of the data file (for format detection).
    objective : str
        Campaign objective / user prompt.
    channels : list | None
        Selected channels from frontend (Social, Display, SEM, etc.)

    Returns
    -------
    BytesIO: PPTX file bytes.
    """
    if channels is None:
        channels = []

    # ---- Parse data file (for insight context + charts) ----
    df = pd.DataFrame()
    data_summary_text = "No structured campaign data uploaded."

    if data_bytes and data_filename:
        ext = data_filename.lower()
        buf = BytesIO(data_bytes)
        try:
            if ext.endswith(".csv"):
                df = pd.read_csv(buf)
            elif ext.endswith(".xlsx") or ext.endswith(".xls"):
                df = pd.read_excel(buf)
            else:
                data_summary_text = f"Unsupported file extension: {ext}"
        except Exception as e:
            data_summary_text = f"Error parsing data file: {e}"
        else:
            data_summary_text = _summarise_dataframe(df)

    # ---- LLM: insights ----
    insights_text = _generate_llm_insights(
        objective=objective,
        data_summary=data_summary_text,
        channels=channels,
        has_image=image_bytes is not None,
    )

    # ---- LLM: recommendations ----
    recos_text = _generate_llm_recos(objective, insights_text)

    # ---- Build themed PPT ----
    ppt_buf = build_final_ppt(
        insights_text=insights_text,
        recos_text=recos_text,
        df=df,
        image_bytes=image_bytes,
    )
    return ppt_buf
