import base64
import tempfile
from pathlib import Path

import streamlit as st
import pandas as pd
from io import BytesIO
from PIL import Image
from chromadb.api.fastapi import FastAPI
from crewai import Agent, Task, Crew
from matplotlib import pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import os
import json, os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Union

st.set_page_config(page_title="Post Campaign Analysis Generator", layout="centered")

# ---------------------------
# Define CrewAI Agents
# ---------------------------

visual_analyst = Agent(
    role="Visual Insight Analyst",
    goal=(
        "Interpret campaign dashboard screenshots to extract KPIs, trends, anomalies, "
        "and propose the best visualizations (charts and illustrative images) to explain them."
    ),
    backstory=(
        "You are a performance marketing analytics specialist who reads complex dashboards, "
        "charts, and graphs daily. You excel at turning visual data into clear insights and "
        "you also think like a designer, suggesting compelling visuals that tell the story."
    ),
llm="gpt-4o-mini",
    llm_config={
        "max_output_tokens": 1000,   # limit output to ~500 tokens
        "temperature": 0.4,
    }
)
# visual_task = Task(
#     description="Analyze the uploaded dashboard image and extract performance metrics, KPIs, trends, and anomalies.",
#     expected_output="A detailed text summary of visual insights with key numbers and campaign performance observations.",
#     agent=visual_analyst,
# )

data_analyst = Agent(
    role="Campaign Data Analyst",
    goal=(
        "Analyze campaign CSV/XLSX data to surface performance trends, ROI, CTR, conversion patterns, "
        "and specify which charts/images should be generated to showcase these insights."
    ),
    backstory=(
        "You are a senior data scientist specializing in digital marketing. You know which metrics "
        "matter to stakeholders and how to structure them into visuals and stories that drive decisions."
    ),
llm="gpt-4o-mini",
    llm_config={
        "max_output_tokens": 100,   # limit output to ~500 tokens
        "temperature": 0.4,
    }
)
# data_task = Task(
#     description="Analyze campaign data from the uploaded CSV or Excel file. Identify key metrics, performance trends, and audience insights.",
#     expected_output="Structured text insights summarizing top-performing metrics, trends over time, and any anomalies or correlations.",
#     agent=data_analyst
# )

insight_synthesizer = Agent(
    role="Insight Synthesizer",
    goal=(
        "Merge visual and data insights into a single, prioritized post-campaign narrative, "
        "including which visuals/images should appear on which slides."
    ),
    backstory=(
        "You are a marketing strategist who is excellent at turning raw analytics into a concise story. "
        "You think in terms of sections in a deck: context, what happened, what worked, what didn’t, "
        "and what to do next. You also map insights to impactful visuals."
    ),
llm="gpt-4o-mini",
    llm_config={
        "max_output_tokens": 100,   # limit output to ~500 tokens
        "temperature": 0.4,
    }
)
# synthesis_task = Task(
#     description="Combine both image-based and data-based insights to create a unified post-campaign performance story.",
#     expected_output="A narrative summary highlighting key takeaways, successes, and improvement opportunities.",
#     agent=insight_synthesizer
# )

ppt_writer = Agent(
role="Presentation Designer",
    goal=(
        "Transform the synthesized insights into a PowerPoint-ready structure with slide titles, bullet points, "
        "and precise specifications for which charts and generated images to place on each slide."
    ),
    backstory=(
        "You are an executive-level presentation designer. You know how to present analytics to C-level stakeholders. "
        "You write concise headlines, sharp bullets, and give very clear instructions on visuals "
        "so a downstream system can generate the PPT and images automatically."
    ),
    llm="gpt-4o-mini",
    llm_config={
        "max_output_tokens": 500,   # limit output to ~500 tokens
        "temperature": 0.4,
    }
)
# ppt_task = Task(
#     description="Create a PowerPoint outline for a post-campaign analysis presentation using the synthesized insights.",
#     expected_output="A structured text outline with slide titles, bullet points, and suggestions for visuals or charts.",
#     agent=ppt_writer
# )

# ---------------------------
# Streamlit UI
# ---------------------------

st.title("📊 AI Post Campaign Analysis Generator")
st.write("Upload your **dashboard image** and **campaign CSV/XLSX file** to generate insights and a PowerPoint report.")

uploaded_image = st.file_uploader("📸 Upload Dashboard/Chart Image", type=["jpg", "jpeg", "png"])
uploaded_file = st.file_uploader("📈 Upload Campaign Data File", type=["csv", "xlsx"])

prompt = st.text_area("Enter your prompt:", placeholder="Ask me something...")

def _rgb_from_list_or_tuple(rgb):
    """Accept list/tuple like [0,166,118] or dict {'r':0,...}"""
    if isinstance(rgb, dict):
        return RGBColor(rgb.get("r",0), rgb.get("g",0), rgb.get("b",0))
    if isinstance(rgb, (list, tuple)) and len(rgb) >= 3:
        return RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))
    return RGBColor(0,0,0)  # default black

def create_ppt_from_json(json_input: Union[str, dict]) -> bytes:
    """Create a PPTX from a JSON structure and return it as bytes for Streamlit download."""
    # Parse input
    if isinstance(json_input, dict):
        data = json_input
    elif isinstance(json_input, str):
        if os.path.exists(json_input):
            with open(json_input, "r", encoding="utf-8") as f:
                data = json.load(f)
        else:
            data = json.loads(json_input)
    else:
        raise ValueError("json_input must be dict, JSON string, or filepath")

    prs = Presentation()

    # Theme
    theme = data.get("theme", {})
    title_color = _rgb_from_list_or_tuple(theme.get("title_color", [0,166,118]))
    content_color = _rgb_from_list_or_tuple(theme.get("content_color", [0,0,0]))
    title_size = Pt(theme.get("title_size", 44))
    content_size = Pt(theme.get("content_size", 24))

    slides = data.get("slides", [])
    for slide_def in slides:
        title = slide_def.get("title", "")
        content = slide_def.get("content", "")
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)

        # Title formatting
        slide.shapes.title.text = title
        if slide.shapes.title.has_text_frame:
            p0 = slide.shapes.title.text_frame.paragraphs[0]
            p0.font.bold = True
            p0.font.size = title_size
            p0.font.color.rgb = title_color

        # Content
        content_box = slide.placeholders[1]
        tf = content_box.text_frame

        if isinstance(content, list):
            bullets = content
        else:
            bullets = [p.strip() for p in str(content).split("\n") if p.strip()]

        if len(bullets) == 0:
            tf.clear()
        else:
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = content_size
                p.font.color.rgb = content_color

        # Speaker notes
        notes_text = slide_def.get("speaker_notes")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = str(notes_text)

    # Return PPTX as bytes for Streamlit download
    from io import BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# -----------------------------
# Helper: Create chart images
# -----------------------------
def create_chart_image(chart, out_dir):
    chart_type = chart["type"].lower()
    title = chart["metric"]
    plt.figure(figsize=(6,4))

    if chart_type == "bar":
        labels = [d["label"] for d in chart["data"]]
        values = [d["value"] for d in chart["data"]]
        plt.bar(labels, values, color='skyblue')
        plt.title(title)
    elif chart_type == "pie":
        labels = [d["label"] for d in chart["data"]]
        values = [d["value"] for d in chart["data"]]
        plt.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
        plt.title(title)
    elif chart_type == "boxplot":
        data_dict = chart["data"]
        # generate synthetic data based on percentiles for visualization
        p25, p50, p75 = data_dict["p25"], data_dict["p50"], data_dict["p75"]
        synthetic_data = [p25]*25 + [p50]*50 + [p75]*25
        plt.boxplot(synthetic_data)
        plt.title(title)
    else:
        plt.close()
        return None

    out_path = Path(out_dir) / f"{title.replace(' ','_')}.png"
    plt.tight_layout()
    plt.savefig(out_path)
    plt.close()
    return str(out_path)



def encode_image_to_data_uri(image_file):
    img_bytes = image_file.getvalue()
    base64_img = base64.b64encode(img_bytes).decode("utf-8")
    return f"data:image/jpeg;base64,{base64_img}"

if st.button("Generate Response"):
    if uploaded_image or uploaded_file:
        with st.spinner("Processing your files..."):
            # Read image
            # image = Image.open(uploaded_image)
            # st.image(image, caption="Uploaded Dashboard Image", use_container_width=True)



            # st.dataframe(df.head())

            # ---------------------------
            # CrewAI Multi-Agent Flow
            # ---------------------------

            # crew = Crew(
            #     agents=[visual_analyst, data_analyst, insight_synthesizer, ppt_writer],
            #     tasks=[visual_task, data_task, synthesis_task, ppt_task]
            # )

            visual_output = ""
            if uploaded_image is not None:
                # 1️⃣ Visual Analysis
                st.info("🔍 Extracting visual insights...")
                data_uri = encode_image_to_data_uri(uploaded_image)
                visual_output = visual_analyst.kickoff(f"Analyze this marketing dashboard and describe key metrics, trends, and anomalies.\n{data_uri}")

            # 2️⃣ Data Analysis
            data_output = ""
            if uploaded_file is not None:
                st.info("📊 Extracting data insights...")
                # Read data
                if uploaded_file.name.endswith(".csv"):
                    df = pd.read_csv(uploaded_file)
                else:
                    df = pd.read_excel(uploaded_file)
                data_summary = df.describe(include='all').to_string()
                # data_output = data_analyst.execute({
                #     "data": df
                # })
                # #data_output = data_task.run(input=data_summary)
                data_output = data_analyst.kickoff(f"Analyze the following campaign data and extract key trends, insights, and performance summary:\n{data_summary}")

            # 3️⃣ Merge both
            st.info("🧩 Merging insights...")
            # synth_input = f"Visual insights:\n{visual_output}\n\nData insights:\n{data_output}"
            synth_input = f"Visual insights:\n{visual_output}\n\nData insights:\n{data_output}"
            synth_output = insight_synthesizer.kickoff(f"Combine the following insights into a unified post-campaign story:\n{synth_input}")
            #synth_output = synthesis_task.run(input={"visual": visual_output, "data": data_output})
            # synth_output = synthesis_task.execute({
            #     "visual": visual_output,
            #     "data": data_output
            # })

            # 4️⃣ Create PPT structure
            st.info("🖋️ Creating PowerPoint outline...")

            # ppt_template = [
            #     {"slide_number": 1, "title": "Campaign Objective", "content": [], "charts": []},
            #     {"slide_number": 2, "title": "Strategy & Execution", "content": [], "charts": []},
            #     {"slide_number": 3, "title": "Performance Growth", "content": [], "charts": []},
            #     {"slide_number": 4, "title": "Performance by Channels", "content": [], "charts": []}
            # ]

            #ppt_input = "create a json format for the below pages generally the json will be used to create the ppt Page 1 Title Page Post Campaign Analysis Page 2  Page Title Campaign  Objective Page 3  Page Title Strategy & Excecutive Page 4  Page Title Performance by Channels Page 5  Page Title Learning and Insights Page 3 and page 4 must have a image location to attach a image I should be able to control the font, font colour of title and description So save the json format in the variable or function then  call the json to convert to pptx use the pptx library in the code give option to customise the style of the ppt"


#             ppt_input = f"""
# Slide Template: {ppt_template}
# Synthesized Insights: {synth_output}
#
# Task:
# 1. From the synthesized insights, generate a structured KPI dataset for the campaign.
#    - Include all relevant metrics like CTR, CPC, Conversions, ROAS, Engagement Rate, etc.
#    - Organize KPIs with date/time or segment breakdowns if available.
# 2. Map the insights into the 4-slide template:
#    Slide 1: Campaign Objective
#    Slide 2: Strategy & Execution
#    Slide 3: Performance Growth
#    Slide 4: Performance by Channels
# 3. For each slide:
#    - Provide concise bullet points summarizing the insights.
#    - Specify charts with actual KPI values so they can be plotted automatically.
# 4. Output must be in JSON format with the following structure:
#    - slide_number
#    - title
#    - bullets
#    - charts: each chart contains type, metric, and data (KPI values)
# 5. Keep output concise, under 500 tokens.
#
# Generate slide-ready proper parseable JSON format  in one output.
# """

            ppt_outline = ppt_writer.kickoff(f"Create 5-7 PowerPoint slide outlines (title + 3-5 bullet points each) for the following insights:\n{synth_output}")
            # ppt_outline = ppt_task.run(input=synth_output)
            # ppt_outline = ppt_task.execute({
            #     "insights": synth_output
            # })

            # ppt_outline = ppt_writer.kickoff(str(ppt_input))
            #
            # print(ppt_outline)  # see what it prints
            # print(ppt_outline.raw)  # raw string (might be empty)
            # print(ppt_outline.json_dict)

            # ---------------------------
            # PPT Generation
            # ---------------------------

            def create_ppt_from_text(ppt_outline_text):
                prs = Presentation()
                slides = ppt_outline_text.split("\n\n")
                for section in slides:
                    if not section.strip():
                        continue
                    lines = section.strip().split("\n")
                    title = lines[0][:80]  # first line as title
                    bullets = lines[1:]
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = title
                    if len(slide.placeholders) > 1:
                        slide.placeholders[1].text = "\n".join(bullets)
                output = BytesIO()
                prs.save(output)
                output.seek(0)
                return output


            ppt_outline_text = getattr(ppt_outline, "raw", None) or str(ppt_outline)

            ppt_bytes = create_ppt_from_text(ppt_outline_text)






            #ppt_outline_text = getattr(ppt_outline, "raw", None) or str(ppt_outline)

            # json_output = json.loads(ppt_outline.raw)

            #print(json_output)



            # # -----------------------------
            # # Generate PPT
            # # -----------------------------
            # prs = Presentation()
            # tmp_dir = Path(tempfile.mkdtemp()) / "charts"
            # tmp_dir.mkdir(parents=True, exist_ok=True)
            #
            # for slide_data in ppt_outline:
            #     slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
            #     slide.shapes.title.text = slide_data["title"]
            #
            #     # Add bullets
            #     tf = slide.placeholders[1].text_frame
            #     tf.clear()
            #     for bullet in slide_data.get("bullets", []):
            #         p = tf.add_paragraph()
            #         p.text = bullet
            #         p.level = 0
            #         p.font.size = Pt(16)
            #
            #     # Add charts
            #     for chart in slide_data.get("charts", []):
            #         chart_path = create_chart_image(chart, tmp_dir)
            #         if chart_path:
            #             slide.shapes.add_picture(chart_path, Inches(1), Inches(2), width=Inches(8), height=Inches(4))
            #
            # # -----------------------------
            # # Save PPT
            # # -----------------------------
            # output_file = "campaign_ppt_with_charts.pptx"
            # prs.save(output_file)


            # ---------------------------
            # Output Download
            # ---------------------------

            st.success("✅ Post Campaign Analysis PPT generated!")
            st.download_button(
                label="📥 Download PowerPoint Report",
                data=ppt_bytes,
                file_name="Post_Campaign_Analysis.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.warning("Please upload either Image or Data files to continue.")
