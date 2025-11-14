# crew_json_to_ppt.py
# -------------------
# Single-file simulation of a 3-agent crew that produces a 13-slide
# post-campaign PPT from Agent inputs (strings) via a validated JSON schema.
#
# Usage:
#   1) pip install python-pptx
#   2) python crew_json_to_ppt.py
#
# Output:
#   - slides.json      (the generated JSON for the 13 slides)
#   - Fully_Fit_PostCampaign.pptx

import json
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Union, List

# -----------------------
# JSON -> PPT helper
# -----------------------
def _rgb_from_list_or_tuple(rgb):
    """Accept list/tuple like [0,166,118] or dict {'r':0,...}"""
    if isinstance(rgb, dict):
        return RGBColor(rgb.get("r",0), rgb.get("g",0), rgb.get("b",0))
    if isinstance(rgb, (list, tuple)) and len(rgb) >= 3:
        return RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))
    # fallback black
    return RGBColor(0,0,0)

def create_ppt_from_json(json_input: Union[str, dict], output_path: str = "Presentation.pptx"):
    """
    Create a PPTX from a JSON structure.
    Accepts dict, JSON string, or path to .json file.
    Schema accepted (partial):
      { "theme": {...}, "slides": [ { "title":..., "content":..., "speaker_notes": "...", "visuals": [...] }, ... ] }
    This function sets title/style, slide text, and speaker notes.
    """
    # Parse input
    if isinstance(json_input, dict):
        data = json_input
    elif isinstance(json_input, str):
        if os.path.exists(json_input):
            with open(json_input, "r", encoding="utf-8") as f:
                data = json.load(f)
        else:
            data = json.loads(json_input)
    else:
        raise ValueError("json_input must be dict, JSON string, or filepath")

    prs = Presentation()

    # Theme
    theme = data.get("theme", {})
    title_color = _rgb_from_list_or_tuple(theme.get("title_color", [0,166,118]))
    content_color = _rgb_from_list_or_tuple(theme.get("content_color", [0,0,0]))
    title_size = Pt(theme.get("title_size", 44))
    content_size = Pt(theme.get("content_size", 24))

    slides = data.get("slides", [])
    for slide_def in slides:
        title = slide_def.get("title", "")
        content = slide_def.get("content", "")
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)

        # Title formatting
        slide.shapes.title.text = title
        if slide.shapes.title.has_text_frame:
            tframe = slide.shapes.title.text_frame
            p0 = tframe.paragraphs[0]
            p0.font.bold = True
            p0.font.size = title_size
            p0.font.color.rgb = title_color

        # Content (string or list)
        content_box = slide.placeholders[1]
        tf = content_box.text_frame

        if isinstance(content, list):
            paras = [str(p) for p in content]
        else:
            paras = [p.strip() for p in str(content).split("\n\n") if p.strip()]

        if len(paras) == 0:
            tf.clear()
        else:
            tf.text = paras[0]
            for extra in paras[1:]:
                new_p = tf.add_paragraph()
                new_p.text = extra

        for paragraph in tf.paragraphs:
            paragraph.font.size = content_size
            paragraph.font.color.rgb = content_color

        # Speaker notes (if provided)
        notes_text = slide_def.get("speaker_notes")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = str(notes_text)

    prs.save(output_path)
    print(f"Presentation saved as {output_path}")

# -----------------------
# Agent 1 and Agent 2 (inputs)
# -----------------------
AGENT_1 = (
    "The campaign achieved strong reach and impressions, with a 22% lift compared to the previous period. "
    "Engagement peaked in Week 2, driven by video creatives that delivered the highest click-through rates."
)

AGENT_2 = (
    "Audience segments aged 18–24 showed the best conversion efficiency, contributing 40% of total actions "
    "despite representing only 28% of impressions. Retargeting ads outperformed prospecting, indicating high "
    "intent among previously exposed users."
)


# -----------------------
# Agent 3: Insights -> JSON converter (function)
# -----------------------
def json_to_ppt_agent(insights_list: List[str], user_instructions: str = None) -> dict:
    """
    Validate, normalize, and convert insights into a 13-slide JSON payload.
    Returns a dict with 'theme' and 'slides' (13 slides).
    Each slide includes provenance and a confidence score.
    """
    # Basic normalization / simple extraction (since we have free text)
    # In production you'd call an LLM or a set of parsers to extract numeric KPIs, weeks, audiences, etc.
    # Here we'll do deterministic parsing using keyword searches to populate the deck.

    # Helper to register provenance for each slide
    def provenance_for(*sources):
        return [{"source": s, "type": "agent_input"} for s in sources]

    # Derive some simple KPIs from the inputs (simulated)
    kpis = {
        "reach_lift_pct": None,
        "peak_week": None,
        "top_creative": None,
        "best_segment": None,
        "segment_actions_pct": None,
        "segment_impressions_pct": None,
        "retargeting_vs_prospecting": None
    }

    # parse Agent 1
    a1 = " ".join(insights_list) if isinstance(insights_list, list) else str(insights_list)
    if "22%" in a1 or "22 percent" in a1 or "22 percent" in a1.lower():
        kpis["reach_lift_pct"] = 22
    if "Week 2" in a1 or "Week 2" in a1:
        kpis["peak_week"] = "Week 2"
    if "video" in a1.lower():
        kpis["top_creative"] = "Video creatives (highest CTR)"

    # parse Agent 2
    a2 = " ".join(insights_list) if isinstance(insights_list, list) else str(insights_list)
    if "18" in a2 and "24" in a2:
        kpis["best_segment"] = "Age 18–24"
    if "40%" in a2 or "40 percent" in a2:
        kpis["segment_actions_pct"] = 40
    if "28%" in a2 or "28 percent" in a2:
        kpis["segment_impressions_pct"] = 28
    if "retargeting" in a2.lower() and "prospecting" in a2.lower():
        kpis["retargeting_vs_prospecting"] = "Retargeting outperformed prospecting — higher intent"

    # Build 13-slide deck (titles + content + speaker notes)
    slides = []

    # 1: Cover
    slides.append({
        "slide_number": 1,
        "title": "Post-Campaign Analysis — Campaign X",
        "content": "Post-campaign summary and recommendations\n\nPresented by: Campaign Analytics Team",
        "speaker_notes": "Cover slide. Include client logo if available.",
        "visuals": [],
        "provenance": provenance_for("agent1", "agent2"),
        "confidence": 0.9
    })

    # 2: Executive Summary
    exec_summary = []
    if kpis["reach_lift_pct"]:
        exec_summary.append(f"Reach & Impressions: {kpis['reach_lift_pct']}% lift vs previous period.")
    if kpis["peak_week"]:
        exec_summary.append(f"Engagement peaked in {kpis['peak_week']}.")
    if kpis["best_segment"]:
        exec_summary.append(f"Top performing audience: {kpis['best_segment']} (high conversion efficiency).")
    exec_summary.append("Retargeting delivered stronger performance than prospecting — strong intent signals.")
    slides.append({
        "slide_number": 2,
        "title": "Executive Summary",
        "content": exec_summary,
        "speaker_notes": "High-level bullets for the leadership team. Mention sample size and timeframe.",
        "visuals": [{"type": "kpi_panel", "data_ref": {"reach_lift_pct": kpis["reach_lift_pct"]}}],
        "provenance": provenance_for("agent1", "agent2"),
        "confidence": 0.88
    })

    # 3: Campaign Objectives & KPIs
    slides.append({
        "slide_number": 3,
        "title": "Campaign Objectives & KPIs",
        "content": [
            "Objective: Drive reach, engagement, and conversions.",
            "Primary KPIs: Impressions, Reach, CTR, Conversions, Conversion Efficiency."
        ],
        "speaker_notes": "List the campaign goals and measurement plan.",
        "visuals": [],
        "provenance": provenance_for("agent1"),
        "confidence": 0.8
    })

    # 4: Reach & Impressions
    reach_text = f"Reach increased by {kpis['reach_lift_pct']}% vs prior period." if kpis["reach_lift_pct"] else "Reach improved vs prior period."
    slides.append({
        "slide_number": 4,
        "title": "Reach & Impressions",
        "content": [reach_text, "Top of funnel activity performed well; consider sustaining frequency."],
        "speaker_notes": "Include chart: impressions over time (week-by-week).",
        "visuals": [{"type": "timeseries", "data_ref": "impressions_over_time"}],
        "provenance": provenance_for("agent1"),
        "confidence": 0.85
    })

    # 5: Engagement & Creative Performance
    creative_text = f"Engagement peaked in {kpis['peak_week']} due to {kpis['top_creative']}." if kpis["peak_week"] and kpis["top_creative"] else "Engagement peaked during campaign midpoint; creative iterated to drive CTR."
    slides.append({
        "slide_number": 5,
        "title": "Engagement & Creative Performance",
        "content": [creative_text, "Recommendation: scale top-performing video formats; A/B test short vs long video cuts."],
        "speaker_notes": "Show CTR by creative. Call out top video variants.",
        "visuals": [{"type": "bar_chart", "data_ref": "ctr_by_creative"}],
        "provenance": provenance_for("agent1"),
        "confidence": 0.9
    })

    # 6: Audience Breakdown
    audience_text = (f"Age 18–24: {kpis['segment_actions_pct']}% of actions vs {kpis['segment_impressions_pct']}% of impressions."
                     if kpis["segment_actions_pct"] and kpis["segment_impressions_pct"]
                     else "Segment-level performance: younger cohorts showed strong efficiency.")
    slides.append({
        "slide_number": 6,
        "title": "Audience Breakdown",
        "content": [audience_text, "Consider segment-specific creatives and offers for 18–24 cohort."],
        "speaker_notes": "Include pie chart: actions by age group and impressions by age group.",
        "visuals": [{"type": "pie_chart", "data_ref": "actions_by_age"}],
        "provenance": provenance_for("agent2"),
        "confidence": 0.92
    })

    # 7: Conversion Efficiency & Funnel
    slides.append({
        "slide_number": 7,
        "title": "Conversion Efficiency & Funnel",
        "content": [
            "Conversion efficiency high among 18–24 cohort.",
            "Recommendation: optimize landing pages for mobile and fast conversion flows."
        ],
        "speaker_notes": "Show funnel: impressions → clicks → actions. Highlight drop-off points.",
        "visuals": [{"type": "funnel_chart", "data_ref": "funnel_metrics"}],
        "provenance": provenance_for("agent2"),
        "confidence": 0.9
    })

    # 8: Retargeting vs Prospecting
    rt_text = kpis["retargeting_vs_prospecting"] or "Retargeting performed better than prospecting."
    slides.append({
        "slide_number": 8,
        "title": "Retargeting vs Prospecting",
        "content": [rt_text, "Recommendation: increase retargeting budget; refresh creative cadence for prospecting."],
        "speaker_notes": "Present comparative CPA/ROAS numbers if available.",
        "visuals": [{"type": "compare_table", "data_ref": "rt_vs_prospecting"}],
        "provenance": provenance_for("agent2"),
        "confidence": 0.9
    })

    # 9: Week-by-Week Performance
    slides.append({
        "slide_number": 9,
        "title": "Week-by-Week Performance",
        "content": ["Week 1: build reach", "Week 2: peak engagement (video-led)", "Week 3: drive conversions", "Week 4: retarget & close"],
        "speaker_notes": "Include line charts for CTR and conversions by week; annotate Week 2 spikes.",
        "visuals": [{"type": "timeseries", "data_ref": "weekly_ctr_conversions"}],
        "provenance": provenance_for("agent1", "agent2"),
        "confidence": 0.86
    })

    # 10: Creative Learnings & Recommendations
    slides.append({
        "slide_number": 10,
        "title": "Creative Learnings & Recommendations",
        "content": [
            "Video creatives drove CTR — prioritize short-form video tests.",
            "Personalized creatives for 18–24 cohort improved CVR — scale for acquisition."
        ],
        "speaker_notes": "List creative tests to run next cycle (A/B ideas, durations, CTAs).",
        "visuals": [],
        "provenance": provenance_for("agent1", "agent2"),
        "confidence": 0.88
    })

    # 11: Budget & Media Mix Recommendations
    slides.append({
        "slide_number": 11,
        "title": "Budget & Media Mix Recommendations",
        "content": [
            "Shift +15% budget to retargeting.",
            "Expand video budgets during Weeks 1–3 to capitalize on discovery and engagement.",
            "Hold prospecting budgets until creative refresh is in place."
        ],
        "speaker_notes": "Show sample allocation table and expected lift scenarios.",
        "visuals": [{"type": "table", "data_ref": "budget_scenarios"}],
        "provenance": provenance_for("agent2"),
        "confidence": 0.8
    })

    # 12: Next Steps & Experiment Plan
    slides.append({
        "slide_number": 12,
        "title": "Next Steps & Experiment Plan",
        "content": [
            "1) Scale top video creatives (2 weeks).",
            "2) Launch landing page optimization for mobile (3 weeks).",
            "3) Increase retargeting budget and creative cadence (ongoing)."
        ],
        "speaker_notes": "Assign owners and timelines for experiments. Include A/B test hypotheses.",
        "visuals": [],
        "provenance": provenance_for("agent1", "agent2"),
        "confidence": 0.85
    })

    # 13: Appendix & Data Notes
    slides.append({
        "slide_number": 13,
        "title": "Appendix & Data Notes",
        "content": [
            "Data sources: ad platform reports, analytics, internal tagging.",
            "Assumptions: timeframe = campaign run dates; sample size = N (replace with actual).",
            "Contact: Campaign Analytics Team"
        ],
        "speaker_notes": "Include definitions for metrics and table of raw numbers if requested.",
        "visuals": [],
        "provenance": provenance_for("agent1", "agent2"),
        "confidence": 0.9
    })

    # Compose final JSON payload
    payload = {
        "schema_version": "1.0",
        "theme": {
            "title_color": [0, 166, 118],
            "content_color": [0, 0, 0],
            "title_size": 36,
            "content_size": 20
        },
        "meta": {
            "generated_by": "json_to_ppt_agent",
            "source_agents": ["agent1", "agent2"],
            "user_instructions": user_instructions or ""
        },
        "slides": slides
    }

    # Basic validation: must be exactly 13 slides
    if len(slides) != 13:
        raise ValueError("Agent must output exactly 13 slides")

    return payload

# -----------------------
# Crew orchestrator
# -----------------------
from google.colab import files

def run_crew(
    save_json_path: str = "/content/slides.json",
    output_pptx: str = "/content/Fully_Fit_PostCampaign.pptx"
):

    # Step 1: Collect agent outputs
    insights = [AGENT_1, AGENT_2]
    print("Agent 1 insight:", AGENT_1)
    print("Agent 2 insight:", AGENT_2)

    # Step 2: Agent 3 produces JSON for PPT
    print("Agent 3 (json_to_ppt) generating 13-slide payload...")
    slides_payload = json_to_ppt_agent(
        insights,
        user_instructions="Create a 13-slide post-campaign deck."
    )

    # Save JSON artifact
    with open(save_json_path, "w", encoding="utf-8") as f:
        json.dump(slides_payload, f, indent=2, ensure_ascii=False)

    print(f"Saved slides JSON to: {save_json_path}")

    # Step 3: Build PPTX
    create_ppt_from_json(slides_payload, output_path=output_pptx)
    print(f"PPTX created at: {output_pptx}")

    # Step 4: Download both files to local machine
    print("Downloading JSON...")
    files.download(save_json_path)

    print("Downloading PPTX...")
    files.download(output_pptx)

    print("Downloads complete!")


# -----------------------
# CLI entrypoint
# -----------------------
if __name__ == "__main__":
    run_crew()
